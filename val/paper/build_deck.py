#!/usr/bin/env python3
"""Assemble the analysis deck: one figure per slide, with what it shows.

IMAGE-FIRST. Each slide is a figure, a noun-phrase title, and the few numbers
that figure is evidence for. The reasoning belongs in the manuscript; a slide
that repeats it competes with the plot for the same two seconds.

The figures and their captions are read off disk — whatever the analyses
actually produced — so re-running an analysis and re-running this is enough to
refresh the deck. Nothing is transcribed by hand.
"""
from __future__ import annotations

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
FIGDIR = os.path.join(HERE, "figures")
RESDIR = os.path.join(HERE, "results")

# 16:9, the only aspect a projector and a laptop agree on.
W, H = Inches(13.333), Inches(7.5)

INK = RGBColor(0x1C, 0x22, 0x28)
MUTED = RGBColor(0x5A, 0x66, 0x72)
CRIMSON = RGBColor(0xC0, 0x39, 0x2B)
RULE = RGBColor(0xD8, 0xDE, 0xE4)
FONT = "Open Sans"          # falls back to the system sans if absent
MONO = "Menlo"

# Slide order, and the noun-phrase title each figure carries. Titles name the
# THING SHOWN, not a claim about it — a slide title that argues gets read as
# the conclusion before the audience has seen the evidence for it.
SLIDES = [
    ("fig01_sample", "Sample composition across two stands"),
    ("fig02_accuracy", "Measured against reference diameter and height"),
    ("fig03_bland_altman", "Agreement with the reference across the size range"),
    ("fig04_crossplatform", "Concordance between the two handsets"),
    ("fig13_shared_error", "Shared and independent components of measurement error"),
    ("fig06_by_size", "Error by diameter and height class"),
    ("fig07_by_site", "Difference between stands, confounded with tree size"),
    ("fig05_model_diagnostics", "Mixed-model residual diagnostics"),
    ("fig08_sigma", "Reported uncertainty against realised error"),
    ("fig09_budget", "Decomposition of the error budget"),
    ("fig10_equivalence", "Equivalence against operational tolerances"),
    ("fig11_estimator", "Chord and cylinder-tangent inversions compared"),
    ("fig12_impact", "Consequences for diameter class and plot totals"),
    ("fig14_scope", "What this design supports, and what it does not"),
]

# Slides whose result rests on hypothesis tests rather than estimates.
# Multiplicity is not controlled across the result set, so these carry the word
# on the slide itself — a caption a reader skips is not where that belongs.
EXPLORATORY = {"fig06_by_size", "fig07_by_site", "fig05_model_diagnostics"}


def p_value_count() -> int:
    """How many hypothesis tests the result tables actually contain.

    Counted, not asserted. The deck said "~634" until a new table added 24 more
    and the slide went quietly stale — the same failure mode as a hard-coded
    caption, in the one number whose whole purpose is to be honest about how
    many chances the study took.
    """
    import csv as _csv
    n = 0
    for name in sorted(os.listdir(RESDIR)) if os.path.isdir(RESDIR) else []:
        if not name.endswith(".csv"):
            continue
        try:
            with open(os.path.join(RESDIR, name)) as fh:
                rows = list(_csv.reader(fh))
        except OSError:
            continue
        if not rows:
            continue
        cols = [i for i, h in enumerate(rows[0])
                if h.strip().lower() in ("p", "p_value", "pvalue")
                or h.strip().lower().startswith("p_")
                or h.strip().lower().endswith("_p")]
        for row in rows[1:]:
            n += sum(1 for i in cols if i < len(row) and row[i].strip())
    return n


def txbox(slide, x, y, w, h, text, size=18, bold=False, color=INK,
          align=PP_ALIGN.LEFT, font=FONT, spacing=1.0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return box


def rule(slide, x, y, w, color=RULE, thickness=Pt(1.25)):
    from pptx.enum.shapes import MSO_SHAPE
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, thickness)
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def read(path, default=""):
    try:
        with open(path) as fh:
            return fh.read().strip()
    except OSError:
        return default


def image_slot(slide, img, x, y, max_w, max_h):
    """Place `img` centred in the box, scaled to fit without distortion."""
    from PIL import Image
    try:
        iw, ih = Image.open(img).size
    except Exception:
        slide.shapes.add_picture(img, x, y, width=max_w)
        return
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(img, x + int((max_w - w) / 2),
                             y + int((max_h - h) / 2), width=w, height=h)


def figure_slide(prs, name, title, index):
    png = os.path.join(FIGDIR, f"{name}.png")
    if not os.path.exists(png):
        return False
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txbox(slide, Inches(0.65), Inches(0.42), Inches(9.6), Inches(0.7),
          title, size=26, bold=True)
    rule(slide, Inches(0.65), Inches(1.12), Inches(1.1), CRIMSON, Pt(3))
    if name in EXPLORATORY:
        txbox(slide, Inches(10.35), Inches(0.52), Inches(2.4), Inches(0.4),
              "EXPLORATORY", size=11, bold=True, color=CRIMSON,
              align=PP_ALIGN.RIGHT)
        txbox(slide, Inches(9.35), Inches(0.86), Inches(3.4), Inches(0.4),
              f"multiplicity uncontrolled across {P_TESTS} tests",
              size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)

    # Figure on the left, the numbers it is evidence for on the right.
    numbers = read(os.path.join(FIGDIR, f"{name}.numbers.txt"))
    if numbers:
        image_slot(slide, png, Inches(0.65), Inches(1.45),
                   Inches(8.5), Inches(5.3))
        txbox(slide, Inches(9.45), Inches(1.55), Inches(3.3), Inches(5.0),
              numbers, size=12, color=MUTED, font=MONO, spacing=1.25)
    else:
        image_slot(slide, png, Inches(0.65), Inches(1.45),
                   Inches(12.0), Inches(5.3))

    caption = read(os.path.join(FIGDIR, f"{name}.txt"))
    if caption:
        caption = re.sub(r"\s+", " ", caption)
        if len(caption) > 300:
            caption = caption[:297].rsplit(" ", 1)[0] + "…"
        txbox(slide, Inches(0.65), Inches(6.85), Inches(11.4), Inches(0.5),
              f"Figure {index}. {caption}", size=10, color=MUTED)
    return True


def title_slide(prs, headline_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txbox(slide, Inches(0.9), Inches(1.55), Inches(11.5), Inches(1.4),
          "Smartphone measurement of tree diameter and height", size=40, bold=True)
    # THE SUBTITLE CARRIES THE SCOPE, not the room's goodwill. "Validation" on
    # its own invites the reader to hear a claim about the app; naming the two
    # stands and the reference makes the claim exactly as large as the design.
    txbox(slide, Inches(0.9), Inches(2.8), Inches(11.5), Inches(1.0),
          "Agreement with a diameter tape and a 3-point laser on 100 stems\n"
          "across two stands, measured by an iPhone and an Android handset",
          size=19, color=MUTED, spacing=1.3)
    rule(slide, Inches(0.9), Inches(4.05), Inches(1.6), CRIMSON, Pt(4))
    if headline_lines:
        txbox(slide, Inches(0.9), Inches(4.45), Inches(11.5), Inches(2.1),
              "\n".join(headline_lines), size=13, color=INK, spacing=1.45)
    txbox(slide, Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.4),
          "Results describe these two stands. Height results are agreement with "
          "the laser, which shares the app's tangent geometry.",
          size=10, color=MUTED)


def table_slide(prs, csv_name, title, max_rows=14, max_cols=8):
    """One results table, rendered as a real table so numbers stay selectable."""
    import csv as _csv
    path = os.path.join(RESDIR, f"{csv_name}.csv")
    if not os.path.exists(path):
        return False
    with open(path) as fh:
        rows = list(_csv.reader(fh))
    if not rows:
        return False
    header, body = rows[0][:max_cols], [r[:max_cols] for r in rows[1:]]
    truncated_rows = len(body) > max_rows
    truncated_cols = len(rows[0]) > max_cols
    body = body[:max_rows]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txbox(slide, Inches(0.65), Inches(0.42), Inches(11.5), Inches(0.7),
          title, size=26, bold=True)
    rule(slide, Inches(0.65), Inches(1.12), Inches(1.1), CRIMSON, Pt(3))

    shape = slide.shapes.add_table(len(body) + 1, len(header),
                                   Inches(0.65), Inches(1.5),
                                   Inches(12.0), Inches(0.32) * (len(body) + 1))
    table = shape.table
    for c, name in enumerate(header):
        cell = table.cell(0, c)
        cell.text = str(name).replace("_", " ")
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(10); run.font.bold = True
                run.font.name = FONT; run.font.color.rgb = INK
    for r, row in enumerate(body, start=1):
        for c in range(len(header)):
            cell = table.cell(r, c)
            cell.text = str(row[c]) if c < len(row) else ""
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(9)
                    run.font.name = MONO if c else FONT
                    run.font.color.rgb = INK

    notes = []
    if truncated_rows:
        notes.append(f"showing {max_rows} of {len(rows) - 1} rows")
    if truncated_cols:
        notes.append(f"{len(rows[0]) - max_cols} further columns")
    notes.append(f"full table: results/{csv_name}.csv")
    txbox(slide, Inches(0.65), Inches(6.95), Inches(11.5), Inches(0.4),
          " · ".join(notes), size=10, color=MUTED)
    return True


P_TESTS = 0


def build(version=2):
    global P_TESTS
    P_TESTS = p_value_count()
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    headlines = []
    for name, _ in SLIDES:
        h = read(os.path.join(FIGDIR, f"{name}.headline.txt"))
        if h:
            headlines.append("— " + re.sub(r"\s+", " ", h))
    title_slide(prs, headlines[:6])

    n = 0
    for name, title in SLIDES:
        if figure_slide(prs, name, title, n + 1):
            n += 1

    for csv_name, title in [
        ("t01_sample", "Sample summary"),
        ("t02_accuracy", "Agreement with the reference, by measurand and stand"),
        ("t05_anova", "Repeated-measures model results (exploratory)"),
        ("t10_equivalence", "Equivalence tests"),
        ("t13_shared_error", "Shared and independent error components"),
    ]:
        table_slide(prs, csv_name, title)

    out = os.path.join(HERE, f"ForestiX_validation_analysis_v{version}.pptx")
    prs.save(out)
    print(f"{out}\n  {len(prs.slides.__iter__.__self__._sldIdLst)} slides, "
          f"{n} figures")
    return out


if __name__ == "__main__":
    build()
